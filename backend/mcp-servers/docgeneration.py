import os
import re
import glob
import pandas as pd
from io import StringIO
from typing import List
from docx import Document
from pptx import Presentation
from tempfile import gettempdir
from datetime import datetime
import hashlib
# from fastapi import FastAPI
from dotenv import load_dotenv
from mcp.server.fastmcp import FastMCP
from typing import Dict, List
from pydantic import BaseModel
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib import colors
import random
import string
from datetime import datetime, timedelta


load_dotenv()

GEN_DOC_PORT = int(os.getenv("PORT", os.getenv("GEN_DOC_PORT", 8000)))
mcp = FastMCP("GenerateDocuments", port=GEN_DOC_PORT)
BASE_URL = os.path.dirname(os.path.abspath(__file__))
BACKEND_DIR = os.path.abspath(os.path.join(BASE_URL, ".."))
PPT_TEMPLATE_DATA_DIR = os.path.join(BACKEND_DIR, "data", "doctemplates")
ppt_filename = "Orion_Innovation_Template.pptx"
ppt_template_path = os.path.join(PPT_TEMPLATE_DATA_DIR, ppt_filename)
TEMP_DIR = gettempdir()
os.makedirs(TEMP_DIR, exist_ok=True)

# Metadata store to map content hashes to filenames
generated_docs = {}

ORG_NAME = "ORION INNOVATION"
CERT_PREFIX = "OI"
DEFAULT_MARGIN = 40

class Slide(BaseModel):
    layout: int
    placeholders: List[str]

class CertificateData(BaseModel):
    student_name: str
    program_name: str
    issue_date: str
    certificate_id: str = None

# --- Helpers ---

def sanitize_filename(text: str, max_length: int = 50) -> str:
    text = text.replace('\n', ' ').replace('\r', ' ')
    text = re.sub(r'[^\w\s-]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    text = text.replace(' ', '_')
    return text[:max_length] or "document"

def get_content_hash(content: str) -> str:
    return hashlib.sha256(content.encode('utf-8')).hexdigest()

def generate_timestamped_filename(base_name: str, ext: str) -> str:
    timestamp = datetime.utcnow().strftime("%Y%m%d%H%M%S")
    return f"{base_name}_{timestamp}.{ext}"

def combine_placeholders(slides: List[Slide]) -> List[Dict]:
    structured = []
    for slide in slides:
        title = slide.placeholders[0] if slide.placeholders else ""
        content = "\n".join(slide.placeholders[1:]) if len(slide.placeholders) > 1 else ""
        structured.append({
            "layout": slide.layout,
            "placeholders": [title, content] if content else [title]
        })
    return structured

def generate_pdf_filename(base_name: str) -> str:
    timestamp = datetime.utcnow().strftime("%Y%m%d%H%M%S")
    safe = re.sub(r"[^\w\d-]+", "_", base_name)[:30]
    return os.path.join(TEMP_DIR, f"Bonafide_{safe}_{timestamp}.pdf")

def generate_certificate_id() -> str:
    suffix = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))
    return f"{CERT_PREFIX}-{datetime.now().strftime('%Y%m%d')}-{suffix}"


# --- Prompts ---
@mcp.prompt("""
When the user asks to generate a document, choose the appropriate tool:
- Use `Generate_PPT` for presentations. Format slides as: [{"layout": 0, "placeholders": ["Title", "Date"]}, {"layout": 1, "placeholders": ["Title", "Content"]}
- Use `Generate_Word_Doc` for narrative or textual content. Pass the text as a single string.
- Use `Generate_Excel` for tabular data. Accept a CSV string as input.
- Use `Generate_Bonafide_Certificate_PDF` when they need a Bonafide Certificate: 
  pass a JSON with `student_name`, `program_name`, `issue_date`, and optional `certificate_id`.

Always include relevant structure to make the tool call successful.
""")

# --- Tools ---

@mcp.tool(name="Generate_Word_Doc", description="Generate a Word document from input content")
async def generate_word_doc(content: str) -> str:
    content_hash = get_content_hash(content)

    # Check reuse
    if content_hash in generated_docs:
        return generated_docs[content_hash]

    # Generate new file
    doc = Document()
    doc.add_heading("Generated Document", 0)
    doc.add_paragraph(content)

    filename = generate_timestamped_filename(sanitize_filename(content), "docx")
    file_path = os.path.join(TEMP_DIR, filename)
    doc.save(file_path)

    generated_docs[content_hash] = file_path
    return file_path

@mcp.tool(name="Generate_Excel", description="Generate an Excel file from CSV string")
async def generate_excel(csv_data: str) -> str:
    content_hash = get_content_hash(csv_data)
    if content_hash in generated_docs:
        return generated_docs[content_hash]

    df = pd.read_csv(StringIO(csv_data))
    first_line = csv_data.splitlines()[0] if csv_data else "excel_data"
    filename = generate_timestamped_filename(sanitize_filename(first_line), "xlsx")
    file_path = os.path.join(TEMP_DIR, filename)
    df.to_excel(file_path, index=False)

    generated_docs[content_hash] = file_path
    return file_path

@mcp.tool(name="Generate_PPT", description="Generate a PowerPoint presentation from structured slide content")
async def generate_ppt(slides: List[Slide]) -> str:
    print(f"Structured slides received: {slides}")

    # Create a hash of the slide structure for caching
    joined_text = " ".join(" ".join(s.placeholders) for s in slides)
    content_hash = get_content_hash(joined_text)

    if content_hash in generated_docs:
        return generated_docs[content_hash]

    prs = Presentation(ppt_template_path)

    # Iterate over each slide layout
    for layout_idx, slide_layout in enumerate(prs.slide_layouts):
        print(f"Slide Layout {layout_idx}: {slide_layout.name}")
        for ph in slide_layout.placeholders:
            print(f"  Placeholder idx: {ph.placeholder_format.idx}, name: '{ph.name}'")
        print("-" * 40)

    # Combine placeholders into a structured format
    structured_slides = combine_placeholders(slides)
    print(f"Combined placeholders: {structured_slides}")

    for idx, slide_info in enumerate(slides):
        layout_idx = slide_info.layout
        placeholders = slide_info.placeholders

        try:
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            for ph in slide.placeholders:
                try:
                    if ph.placeholder_format.idx == 0 and len(placeholders) > 0:
                        ph.text = placeholders[1]  # heading or title
                    elif ph.placeholder_format.idx == 10 and len(placeholders) > 1:
                        ph.text = placeholders[0]  # content
                except Exception as e:
                    print(f"Warning: Failed to set text in placeholder '{ph.name}' on slide {idx}: {e}")

        except Exception as e:
            print(f"Error adding slide {idx}: {e}")

    filename = generate_timestamped_filename(sanitize_filename(joined_text), "pptx")
    file_path = os.path.join(TEMP_DIR, filename)
    prs.save(file_path)

    generated_docs[content_hash] = file_path
    return file_path

def generate_certificate_id(prefix="ORION") -> str:
    date_part = datetime.now().strftime("%Y%m%d")
    rand_part = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))
    return f"{prefix}-{date_part}-{rand_part}"

@mcp.tool(
    name="Generate_Bonafide_Certificate_PDF",
    description="Generate a Bonafide/Employment Certificate for Orion Innovation."
)
async def generate_certificate_pdf(data: CertificateData) -> str:
    today = datetime.now()
    next_month = (today.replace(day=1) + timedelta(days=32)).replace(day=1) - timedelta(days=1)
    data.issue_date = data.issue_date or today.strftime("%d-%m-%Y")
    data.certificate_id = data.certificate_id or generate_certificate_id()

    key_string = f"{data.student_name}|{data.program_name}|{data.issue_date}|{data.certificate_id}"
    content_hash = hashlib.sha256(key_string.encode("utf-8")).hexdigest()
    if content_hash in generated_docs:
        return generated_docs[content_hash]

    pdf_path = generate_pdf_filename(data.student_name)
    c = canvas.Canvas(pdf_path, pagesize=landscape(A4))
    width, height = landscape(A4)
    margin = 50

    # Borders
    c.setStrokeColor(colors.HexColor("#000000"))
    c.setLineWidth(3)
    c.rect(margin, margin, width - 2 * margin, height - 2 * margin)

    c.setStrokeColor(colors.lightgrey)
    c.setLineWidth(1)
    c.rect(margin + 10, margin + 10, width - 2 * (margin + 10), height - 2 * (margin + 10))

    # Header
    c.setFont("Helvetica-Bold", 22)
    c.setFillColor(colors.HexColor("#000000"))
    c.drawCentredString(width / 2, height - margin - 30, "Orion Innovation")

    c.setFont("Helvetica", 12)
    c.setFillColor(colors.black)
    c.drawCentredString(width / 2, height - margin - 50, "Global Technology Services | www.orioninc.com")

    # Title
    c.setFont("Helvetica-Bold", 26)
    c.setFillColor(colors.black)
    c.drawCentredString(width / 2, height - margin - 100, "BONAFIDE CERTIFICATE")

    # Body Content
    content_text = (
        f"This is to certify that Mr./Ms. {data.student_name} is working as a "
        f"in Orion Innovation. This certificate is issued for his {data.program_name} related purposes.\n\n"
        f"This is valid till: {next_month.strftime('%d-%m-%Y')}"
    )

    c.setFont("Helvetica", 14)
    text_obj = c.beginText()
    text_obj.setTextOrigin(margin + 70, height - margin - 150)
    text_obj.setLeading(22)
    max_width = width - 2 * (margin + 70)

    for line in content_text.split("\n"):
        words = line.split()
        buffer = ""
        for word in words:
            test_line = f"{buffer} {word}".strip()
            if c.stringWidth(test_line, "Helvetica", 14) < max_width:
                buffer = test_line
            else:
                text_obj.textLine(buffer)
                buffer = word
        if buffer:
            text_obj.textLine(buffer)
    c.drawText(text_obj)

    # Signature
    c.setFont("Helvetica", 12)
    c.drawString(width - margin - 220, margin + 60, "HR Manager – Orion Innovation")
    c.drawString(width - margin - 220, margin + 80, "__________________________")

    # Certificate ID
    c.setFont("Helvetica-Oblique", 10)
    c.drawString(margin + 10, margin + 40, f"Certificate ID: {data.certificate_id}")

    # Watermark (bounded within border)
    c.saveState()
    c.setFont("Helvetica-Bold", 40)
    c.setFillGray(0.90, 0.5)
    c.translate(width / 2, height / 2)
    c.rotate(45)
    c.drawCentredString(0, 0, "ORION INNOVATION")
    c.restoreState()
    c.save()

    generated_docs[content_hash] = pdf_path
    return pdf_path


@mcp.tool(name="Cleanup_Temp_Files", description="Delete all generated temporary document files")
async def cleanup_temp_files() -> str:
    deleted = 0
    for ext in ["docx", "pptx", "xlsx"]:
        pattern = os.path.join(TEMP_DIR, f"*.{ext}")
        for filepath in glob.glob(pattern):
            try:
                os.remove(filepath)
                deleted += 1
            except Exception:
                pass
    generated_docs.clear()
    return f"Deleted {deleted} temporary files."

# --- Server Start ---
if __name__ == "__main__":
    mcp.run(transport="sse")
